"""
增强型PPT加载器 (Enhanced PowerPoint Loader)

该模块提供了同步和异步两个版本的PPT加载器，集成了大型视觉语言模型（LLM）
来分析PPT中的文本、图片和图表等内容，并能自动生成问答（QA）对，
适用于构建检索增强生成（RAG）系统。

核心功能：
1. 使用 `unstructured` 提取PPT中的基础文本。
2. 将每张幻灯片渲染为图片，或提取幻灯片中的独立图片。
3. 调用视觉大模型（如GPT-4o）分析幻灯片视觉内容（图表、图片等）。
4. 基于文本和视觉分析结果，为每张幻灯片生成高质量的QA对。
5. 所有生成的内容（文本、分析、QA）都以丰富的元数据封装在Document对象中。
6. 提供同步 (`EnhancedPPTLoader`) 和高并发的异步 (`AsyncEnhancedPPTLoader`) 版本。
"""

import asyncio
import base64
import json
import os
from typing import Dict, List, Tuple, Any

import openai
from langchain_community.document_loaders import UnstructuredPowerPointLoader
from pptx import Presentation

# 假设项目中有一个统一的Logger配置
from logging import getLogger
logger = getLogger(__name__)


class Document:
    """一个简化的数据结构，用于模仿LangChain的Document，用于演示目的。"""
    def __init__(self, page_content: str, metadata: dict):
        self.page_content = page_content
        self.metadata = metadata

    def __repr__(self):
        return f"Document(page_content='{self.page_content[:50]}...', metadata={self.metadata})"


class BaseLoader:
    """加载器的基类，用于定义通用接口和辅助函数。"""
    def validate_file(self, file_path: str) -> bool:
        """检查文件是否存在且非空。"""
        return os.path.exists(file_path) and os.path.getsize(file_path) > 0


class EnhancedPPTLoader(BaseLoader):
    """
    增强型PowerPoint文件加载器。
    它不仅提取文本，还利用视觉大模型分析幻灯片中的图片和图表，
    并自动生成可用于RAG系统的问答对。
    """

    # 默认的Prompts，可以被实例覆盖
    DEFAULT_VISION_PROMPT = (
        "你是一位专业的演示文稿分析师。请详细分析这张幻灯片截图的内容，遵循以下结构："
        "1. **核心主题**：一句话总结这张幻灯片的主要议题。"
        "2. **视觉元素分析**：描述图片、图表、流程图等关键视觉元素。如果是图表，请分析其类型（如柱状图、折线图）、展示的数据、趋势和关键洞见。"
        "3. **关键信息提取**：列出从幻灯片中可以提取出的3-5个关键信息点或数据。"
        "4. **图文关系**：结合下方提供的幻灯片文本，分析视觉内容与文本信息是如何相互支持或补充的。"
        "请使用中文以Markdown格式清晰地回答。"
    )
    
    DEFAULT_QA_PROMPT_TEMPLATE = """
    你是一个出色的问答生成器。基于以下单张PPT幻灯片的文本和视觉分析内容，生成{num_qa_pairs}个高质量的问答对(QA对)。
    要求：
    1. 问题需具体且有深度，能够引导出幻灯片的核心知识。
    2. 答案需准确、详细，完全基于所提供的上下文内容。
    3. 问答对需要全面覆盖文本和视觉分析中的关键信息。
    4. 对于答案主要源于视觉内容（如分析图表得出的结论）的问题，请在答案末尾标注 "[视觉分析]"。
    5. 返回一个JSON格式的列表，不要包含任何额外的解释。

    幻灯片上下文：
    ---
    {slide_context}
    ---

    请严格按照以下JSON格式返回：
    [
        {{
            "question": "问题1",
            "answer": "答案1"
        }},
        {{
            "question": "问题2", 
            "answer": "答案2"
        }}
    ]
    """

    def __init__(
        self,
        vision_model: str = "gpt-4o",
        qa_model: str = "gpt-4o-mini",
        openai_api_key: str = None,
        num_qa_pairs: int = 3,
        max_tokens_vision: int = 1500,
        max_tokens_qa: int = 1000,
        vision_prompt: str = None,
        qa_prompt_template: str = None,
    ):
        """
        初始化增强型PPT加载器。

        Args:
            vision_model: 用于图像分析的视觉大模型名称。
            qa_model: 用于生成QA对的语言模型名称。
            openai_api_key: OpenAI API密钥。若为None，则从环境变量 "OPENAI_API_KEY" 获取。
            num_qa_pairs: 为每张幻灯片生成的QA对数量。
            max_tokens_vision: 视觉分析任务的最大token数。
            max_tokens_qa: QA生成任务的最大token数。
            vision_prompt: 自定义视觉分析的prompt。
            qa_prompt_template: 自定义QA生成的prompt模板。
        
        Raises:
            ValueError: 如果无法找到OpenAI API密钥。
        """
        api_key = openai_api_key or os.getenv("OPENAI_API_KEY")
        if not api_key:
            raise ValueError("OpenAI API key is required. Please provide it or set the OPENAI_API_KEY environment variable.")
        
        self.openai_client = openai.OpenAI(api_key=api_key)
        self.vision_model = vision_model
        self.qa_model = qa_model
        self.num_qa_pairs = num_qa_pairs
        self.max_tokens_vision = max_tokens_vision
        self.max_tokens_qa = max_tokens_qa
        self.vision_prompt = vision_prompt or self.DEFAULT_VISION_PROMPT
        self.qa_prompt_template = qa_prompt_template or self.DEFAULT_QA_PROMPT_TEMPLATE

    def load(self, file_path: str) -> List[Document]:
        """
        同步加载PPT文件，提取文本，分析视觉内容，并生成QA对。

        Args:
            file_path: PPT文件的路径。

        Returns:
            包含文本、视觉分析和QA对的文档列表。
        """
        logger.info(f"[EnhancedPPTLoader] 开始处理PPT文件: {file_path}")
        if not self.validate_file(file_path):
            logger.error(f"[EnhancedPPTLoader] 文件无效或未找到: {file_path}")
            raise FileNotFoundError(f"File not found or is empty: {file_path}")

        try:
            # 1. 加载基础文本内容
            text_documents = self._load_text_content(file_path)

            # 2. 分析幻灯片视觉内容
            vision_documents = self._process_slide_images(file_path, text_documents)

            # 3. 合并所有内容并生成QA对
            all_content_docs = text_documents + vision_documents
            qa_documents = self._generate_qa_pairs(all_content_docs)

            final_documents = text_documents + vision_documents + qa_documents
            logger.info(f"[EnhancedPPTLoader] 处理完成，共生成 {len(final_documents)} 个文档。")
            return final_documents

        except Exception as e:
            logger.error(f"[EnhancedPPTLoader] 处理PPT文件 '{file_path}' 时发生未知错误: {e}", exc_info=True)
            raise

    def _load_text_content(self, file_path: str) -> List[Document]:
        """使用UnstructuredPowerPointLoader加载PPT中的文本内容。"""
        logger.info("正在提取文本内容...")
        loader = UnstructuredPowerPointLoader(file_path, mode="elements")
        documents = loader.load()
        
        # 将文本内容按幻灯片页码聚合
        slide_texts: Dict[int, str] = {}
        for doc in documents:
            slide_num = int(doc.metadata.get("page_number", 0))
            if slide_num > 0:
                if slide_num not in slide_texts:
                    slide_texts[slide_num] = ""
                slide_texts[slide_num] += doc.page_content + "\n\n"

        # 为每张幻灯片的聚合文本创建一个Document
        text_documents = []
        for slide_num, content in slide_texts.items():
            text_documents.append(Document(
                page_content=content.strip(),
                metadata={
                    "source_file": os.path.basename(file_path),
                    "slide_number": slide_num,
                    "content_type": "text",
                },
            ))
        return text_documents

    def _render_slide_as_image(self, slide: Any, slide_index: int) -> Tuple[bytes, str]:
        """
        【理想实现】将单张幻灯片渲染为图片。
        注意：此功能的实现依赖于外部库，如 pywin32 (仅Windows) 或通过
        libreoffice/unoconv 进行跨平台转换。此处为一个概念性占位符。
        """
        # 示例：在Windows上使用pywin32 (需要安装PowerPoint)
        # import win32com.client
        # powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        # pres = powerpoint.Presentations.Open(file_path, WithWindow=False)
        # temp_image_path = f"temp_slide_{slide_index}.png"
        # pres.Slides[slide_index].Export(temp_image_path, "PNG")
        # with open(temp_image_path, "rb") as f:
        #     image_data = f.read()
        # os.remove(temp_image_path)
        # pres.Close()
        # powerpoint.Quit()
        # return image_data, "png"
        raise NotImplementedError("幻灯片渲染功能需要配置特定环境的外部依赖（如pywin32）。")

    def _process_slide_images(self, file_path: str, text_documents: List[Document]) -> List[Document]:
        """
        分析每张幻灯片的视觉内容。
        首选策略是渲染整张幻灯片，如果失败，则回退到提取嵌入的图片。
        """
        logger.info("正在分析幻灯片的视觉内容...")
        prs = Presentation(file_path)
        vision_docs = []
        
        # 创建一个从幻灯片号到文本内容的映射
        slide_text_map = {doc.metadata["slide_number"]: doc.page_content for doc in text_documents}

        for i, slide in enumerate(prs.slides):
            slide_num = i + 1
            try:
                # 策略1：尝试将整张幻灯片渲染为图片进行分析（推荐）
                # 注意：_render_slide_as_image 需要您根据环境自行实现
                # image_data, image_format = self._render_slide_as_image(slide, i)
                # analysis = self._analyze_image_with_vision_model(image_data, slide_text_map.get(slide_num, ""))
                # vision_docs.append(Document(
                #     page_content=analysis,
                #     metadata={
                #         "source_file": os.path.basename(file_path),
                #         "slide_number": slide_num,
                #         "content_type": "slide_visual_analysis",
                #         "analysis_source": "full_slide_render",
                #     },
                # ))
                
                # 策略2：作为备用方案，提取幻灯片中嵌入的图片进行分析
                # 如果您无法实现渲染，可以取消下面的注释来使用此策略
                embedded_images = self._extract_images_from_slide(slide)
                if not embedded_images:
                    continue

                # 此处只分析第一张有意义的图片，避免信息冗余
                image_data, image_format = embedded_images[0]
                analysis = self._analyze_image_with_vision_model(image_data, slide_text_map.get(slide_num, ""))
                vision_docs.append(Document(
                    page_content=analysis,
                    metadata={
                        "source_file": os.path.basename(file_path),
                        "slide_number": slide_num,
                        "content_type": "slide_visual_analysis",
                        "analysis_source": "embedded_image",
                        "image_format": image_format,
                    },
                ))

            except NotImplementedError:
                logger.warning(f"幻灯片 {slide_num} 渲染功能未实现，跳过视觉分析。")
                # 如果您只想用备用策略，请注释掉 raise 语句并取消上方策略2的注释。
            except Exception as e:
                logger.error(f"处理幻灯片 {slide_num} 的视觉内容时出错: {e}", exc_info=True)
                
        return vision_docs

    def _extract_images_from_slide(self, slide: Any) -> List[Tuple[bytes, str]]:
        """从单张幻灯片的rels中提取图片数据和格式。"""
        images = []
        for rel in slide.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    image_part = rel.target_part
                    image_data = image_part.blob
                    image_format = image_part.content_type.split('/')[-1]
                    images.append((image_data, image_format))
                except Exception as e:
                    logger.warning(f"提取图片时出错: {e}")
        return images

    def _analyze_image_with_vision_model(self, image_data: bytes, slide_text_content: str) -> str:
        """使用视觉大模型分析图片内容，并结合文本上下文。"""
        image_base64 = base64.b64encode(image_data).decode("utf-8")
        
        prompt_text = self.vision_prompt
        if slide_text_content:
            prompt_text += f"\n\n附加的幻灯片文本内容参考：\n---\n{slide_text_content}"

        messages = [{
            "role": "user",
            "content": [
                {"type": "text", "text": prompt_text},
                {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{image_base64}"}},
            ],
        }]
        
        try:
            response = self.openai_client.chat.completions.create(
                model=self.vision_model,
                messages=messages,
                max_tokens=self.max_tokens_vision
            )
            return response.choices[0].message.content
        except openai.RateLimitError as e:
            logger.error(f"OpenAI API 速率限制: {e}. 等待后重试...")
            raise
        except openai.APIError as e:
            logger.error(f"OpenAI API 错误: {e}")
            raise

    def _generate_qa_pairs(self, documents: List[Document]) -> List[Document]:
        """基于所有文档内容，为每张幻灯片生成QA对。"""
        logger.info("正在生成QA对...")
        qa_documents = []
        
        # 按幻灯片编号对所有内容（文本+视觉分析）进行分组
        slide_contexts: Dict[int, str] = {}
        for doc in documents:
            slide_num = doc.metadata.get("slide_number")
            if slide_num:
                if slide_num not in slide_contexts:
                    slide_contexts[slide_num] = ""
                slide_contexts[slide_num] += f"## 内容来源: {doc.metadata['content_type']} ##\n{doc.page_content}\n\n"
        
        for slide_num, context in slide_contexts.items():
            try:
                qa_list = self._create_qa_for_slide(context, slide_num)
                qa_documents.extend(qa_list)
            except Exception as e:
                logger.error(f"为幻灯片 {slide_num} 生成QA对失败: {e}", exc_info=True)
                # 即使失败，也继续处理下一张幻灯片
                
        return qa_documents

    def _create_qa_for_slide(self, slide_context: str, slide_num: int) -> List[Document]:
        """为单个幻灯片的聚合内容创建QA对。"""
        prompt = self.qa_prompt_template.format(
            num_qa_pairs=self.num_qa_pairs,
            slide_context=slide_context
        )
        
        try:
            response = self.openai_client.chat.completions.create(
                model=self.qa_model,
                messages=[{"role": "user", "content": prompt}],
                response_format={"type": "json_object"},
                max_tokens=self.max_tokens_qa
            )
            # OpenAI的JSON模式有时会把列表包在一个键里，例如 {"qa_pairs": [...] }
            # 这里做一些兼容性处理
            content = response.choices[0].message.content
            qa_data = json.loads(content)
            if isinstance(qa_data, dict) and len(qa_data) == 1:
                qa_list = next(iter(qa_data.values()))
            else:
                qa_list = qa_data

            if not isinstance(qa_list, list):
                raise ValueError("模型返回的JSON不是一个列表")

        except (json.JSONDecodeError, ValueError) as e:
            logger.error(f"解析为幻灯片 {slide_num} 生成的QA JSON失败: {e}. 尝试备用方案。")
            return self._fallback_qa_generation(slide_context, slide_num)
        
        qa_documents = []
        for i, qa_item in enumerate(qa_list):
            if "question" not in qa_item or "answer" not in qa_item:
                logger.warning(f"幻灯片 {slide_num} 的一个QA对格式不正确，已跳过: {qa_item}")
                continue
                
            qa_doc = Document(
                page_content=f"Q: {qa_item['question']}\nA: {qa_item['answer']}",
                metadata={
                    "slide_number": slide_num,
                    "content_type": "qa_pair",
                    "qa_index": i + 1,
                    "question": qa_item['question'],
                    "answer": qa_item['answer'],
                    "source": "generated_qa",
                },
            )
            qa_documents.append(qa_doc)
            
        return qa_documents

    def _fallback_qa_generation(self, context: str, slide_num: int) -> List[Document]:
        """当主QA生成方案失败时的备用方案，生成一个总结性问题。"""
        logger.info(f"正在为幻灯片 {slide_num} 执行备用QA生成...")
        return [Document(
            page_content=f"Q: 请总结幻灯片 {slide_num} 的核心内容。\nA: 这张幻灯片的主要内容摘要如下：{context[:500]}...",
            metadata={
                "slide_number": slide_num,
                "content_type": "qa_pair",
                "qa_index": 1,
                "source": "fallback_qa_summary",
            },
        )]


class AsyncEnhancedPPTLoader(EnhancedPPTLoader):
    """
    异步版本的增强型PPT加载器。
    利用asyncio并发执行网络密集型任务（如API调用），显著提升处理速度。
    """

    async def aload(self, file_path: str) -> List[Document]:
        """
        异步加载PPT文件，并发处理视觉分析和QA生成。

        Args:
            file_path: PPT文件的路径。

        Returns:
            包含文本、视觉分析和QA对的文档列表。
        """
        logger.info(f"[AsyncEnhancedPPTLoader] 开始异步处理PPT文件: {file_path}")
        if not self.validate_file(file_path):
            logger.error(f"[AsyncEnhancedPPTLoader] 文件无效或未找到: {file_path}")
            raise FileNotFoundError(f"File not found or is empty: {file_path}")

        try:
            # 1. 同步加载文本内容（通常很快，无需异步）
            text_documents = self._load_text_content(file_path)

            # 2. 异步并发分析所有幻灯片的视觉内容
            vision_documents = await self._async_process_all_slides(file_path, text_documents)

            # 3. 异步并发为所有幻灯片生成QA对
            all_content_docs = text_documents + vision_documents
            qa_documents = await self._async_generate_all_qa_pairs(all_content_docs)

            final_documents = text_documents + vision_documents + qa_documents
            logger.info(f"[AsyncEnhancedPPTLoader] 异步处理完成，共生成 {len(final_documents)} 个文档。")
            return final_documents

        except Exception as e:
            logger.error(f"[AsyncEnhancedPPTLoader] 处理PPT文件 '{file_path}' 时发生未知错误: {e}", exc_info=True)
            raise
    
    async def _async_process_all_slides(self, file_path: str, text_documents: List[Document]) -> List[Document]:
        """并发处理所有幻灯片的视觉分析。"""
        logger.info("正在并发分析幻灯片的视觉内容...")
        prs = Presentation(file_path)
        slide_text_map = {doc.metadata["slide_number"]: doc.page_content for doc in text_documents}
        
        tasks = []
        for i, slide in enumerate(prs.slides):
            slide_num = i + 1
            # 在异步环境中，最好不要依赖需要外部进程的渲染（除非该进程也是异步管理的）
            # 因此，这里直接使用图片提取策略
            embedded_images = self._extract_images_from_slide(slide)
            if embedded_images:
                image_data, image_format = embedded_images[0]
                slide_text = slide_text_map.get(slide_num, "")
                task = self._async_analyze_image(image_data, image_format, slide_num, slide_text, os.path.basename(file_path))
                tasks.append(task)
        
        vision_docs = []
        results = await asyncio.gather(*tasks, return_exceptions=True)
        for res in results:
            if isinstance(res, Document):
                vision_docs.append(res)
            elif isinstance(res, Exception):
                logger.error(f"一个异步视觉分析任务失败: {res}")
        return vision_docs

    async def _async_analyze_image(self, image_data: bytes, image_format: str, slide_num: int, slide_text: str, filename: str) -> Document:
        """异步调用视觉模型分析单张图片。"""
        loop = asyncio.get_event_loop()
        # 使用run_in_executor在线程池中运行同步的SDK调用，防止阻塞事件循环
        analysis = await loop.run_in_executor(
            None, self._analyze_image_with_vision_model, image_data, slide_text
        )
        return Document(
            page_content=analysis,
            metadata={
                "source_file": filename,
                "slide_number": slide_num,
                "content_type": "slide_visual_analysis",
                "analysis_source": "embedded_image",
                "image_format": image_format,
            },
        )
        
    async def _async_generate_all_qa_pairs(self, documents: List[Document]) -> List[Document]:
        """并发为所有幻灯片生成QA对。"""
        logger.info("正在并发生成QA对...")
        slide_contexts: Dict[int, str] = {}
        for doc in documents:
            slide_num = doc.metadata.get("slide_number")
            if slide_num:
                if slide_num not in slide_contexts:
                    slide_contexts[slide_num] = ""
                slide_contexts[slide_num] += f"## 内容来源: {doc.metadata['content_type']} ##\n{doc.page_content}\n\n"

        tasks = [self._async_create_qa_for_slide(ctx, num) for num, ctx in slide_contexts.items()]
        
        qa_documents = []
        results = await asyncio.gather(*tasks, return_exceptions=True)
        for res in results:
            if isinstance(res, list):
                qa_documents.extend(res)
            elif isinstance(res, Exception):
                logger.error(f"一个异步QA生成任务失败: {res}")
        return qa_documents

    async def _async_create_qa_for_slide(self, slide_context: str, slide_num: int) -> List[Document]:
        """异步为单个幻灯片创建QA对。"""
        loop = asyncio.get_event_loop()
        # 同样使用线程池运行同步的SDK调用
        return await loop.run_in_executor(
            None, self._create_qa_for_slide, slide_context, slide_num
        )


# # ================== 示例用法 ==================
# async def main():
#     # 配置日志记录
#     import logging
#     logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

#     # 确保已设置 OPENAI_API_KEY 环境变量
#     if "OPENAI_API_KEY" not in os.environ:
#         print("错误：请先设置 OPENAI_API_KEY 环境变量。")
#         return

#     # 替换为你的PPT文件路径
#     ppt_file_path = "path/to/your/presentation.pptx"
#     if not os.path.exists(ppt_file_path):
#         print(f"错误：示例文件 {ppt_file_path} 不存在。请创建一个或修改路径。")
#         # 创建一个虚拟PPT文件用于测试
#         from pptx import Presentation
#         prs = Presentation()
#         slide = prs.slides.add_slide(prs.slide_layouts[5])
#         slide.shapes.title.text = "这是一个测试标题"
#         slide.shapes.add_textbox(100, 200, 500, 100).text = "这是幻灯片的主要内容，讨论了增长策略。"
#         # 您还可以在这里添加图片以测试视觉分析功能
#         prs.save(ppt_file_path)
#         print(f"已创建测试文件：{ppt_file_path}")

#     # --- 同步加载器示例 ---
#     print("\n--- 开始同步加载 ---")
#     sync_loader = EnhancedPPTLoader()
#     try:
#         sync_documents = sync_loader.load(ppt_file_path)
#         print(f"同步加载完成，获取到 {len(sync_documents)} 个文档。")
#         for doc in sync_documents:
#             print(doc)
#     except Exception as e:
#         print(f"同步加载失败: {e}")

#     # --- 异步加载器示例 ---
#     print("\n--- 开始异步加载 ---")
#     async_loader = AsyncEnhancedPPTLoader()
#     try:
#         async_documents = await async_loader.aload(ppt_file_path)
#         print(f"异步加载完成，获取到 {len(async_documents)} 个文档。")
#         for doc in async_documents:
#             print(doc)
#     except Exception as e:
#         print(f"异步加载失败: {e}")

# if __name__ == "__main__":
#     asyncio.run(main())